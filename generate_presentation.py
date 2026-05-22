# -*- coding: utf-8 -*-
"""
Генератор презентации для защиты ВКР.
Тема: Классификация растений по изображениям с использованием
      алгоритмов компьютерного зрения для агротехнического контроля.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DIR = os.path.dirname(os.path.abspath(__file__))

# ── Цветовая схема (чёрно-белая) ──
BG_DARK    = RGBColor(0xFF, 0xFF, 0xFF)   # белый фон
BG_ACCENT  = RGBColor(0xF0, 0xF0, 0xF0)   # светло-серый акцент
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x55, 0x55, 0x55)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
GREEN      = RGBColor(0x00, 0x00, 0x00)
BLUE       = RGBColor(0x33, 0x33, 0x33)
ORANGE     = RGBColor(0x33, 0x33, 0x33)
RED        = RGBColor(0x33, 0x33, 0x33)
YELLOW     = RGBColor(0x33, 0x33, 0x33)


def set_slide_bg(slide, color):
    """Установить цвет фона слайда."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=20,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Добавить текстовое поле."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=BLACK, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    """Добавить параграф в текстовый фрейм."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    """Добавить изображение, если файл существует."""
    path = os.path.join(DIR, img_name)
    if os.path.isfile(path):
        kwargs = {}
        if width:  kwargs['width'] = width
        if height: kwargs['height'] = height
        slide.shapes.add_picture(path, left, top, **kwargs)
        return True
    return False


def add_rounded_rect(slide, left, top, width, height, fill_color, text='',
                     font_size=16, font_color=BLACK):
    """Добавить скруглённый прямоугольник с текстом."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
    return shape


# ═══════════════════════════════════════════════════
#                     СЛАЙДЫ
# ═══════════════════════════════════════════════════

def slide_title(prs):
    """Слайд 1: Титульный."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # ── Шапка университета (изображение) ──
    add_image_safe(slide, 'university_header.png', Cm(0.5), Cm(0.3), width=Cm(24.4))

    # ── Заголовок «ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА» ──
    add_textbox(slide, Cm(2), Cm(5), Cm(21), Cm(1.2),
                'ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА',
                font_size=16, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

    # Тема
    tf = add_textbox(slide, Cm(1.5), Cm(6.8), Cm(22), Cm(4),
                     'Классификация растений по изображениям',
                     font_size=28, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, 'с использованием алгоритмов компьютерного зрения',
                  font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, 'для агротехнического контроля',
                  font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Автор
    tf2 = add_textbox(slide, Cm(2), Cm(12.5), Cm(21), Cm(3),
                      'Выполнил: Иманалиев Асылбек',
                      font_size=16, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, 'Научный руководитель: ___________________',
                  font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Год
    add_textbox(slide, Cm(2), Cm(16.5), Cm(21), Cm(1.5),
                'Бишкек — 2025', font_size=16, color=BLACK, alignment=PP_ALIGN.CENTER)

    # Нижняя линия
    line_bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(18.6), Cm(25.4), Pt(3))
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = BLACK
    line_bottom.line.fill.background()


def slide_problem(prs):
    """Слайд 2: Актуальность и цель."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), BLUE)

    add_textbox(slide, Cm(1), Cm(0.8), Cm(23), Cm(1.5),
                'Актуальность и цель работы',
                font_size=28, color=BLACK, bold=True)

    # Актуальность
    add_rounded_rect(slide, Cm(1), Cm(2.5), Cm(11), Cm(7), BG_ACCENT,
                     '', font_size=12)
    tf = add_textbox(slide, Cm(1.5), Cm(2.7), Cm(10), Cm(1),
                     'Актуальность', font_size=20, color=GREEN, bold=True)
    items = [
        'Потери урожая от болезней растений до 40%',
        'Ручная диагностика медленная и субъективная',
        'Нейросети способны классифицировать с точностью >95%',
        'Потребность в автоматизации агротехнического контроля',
    ]
    for item in items:
        add_paragraph(tf, f'  {item}', font_size=16, color=LIGHT_GRAY)

    # Цель
    add_rounded_rect(slide, Cm(13), Cm(2.5), Cm(11), Cm(7), BG_ACCENT,
                     '', font_size=12)
    tf2 = add_textbox(slide, Cm(13.5), Cm(2.7), Cm(10), Cm(1),
                      'Цель', font_size=20, color=GREEN, bold=True)
    add_paragraph(tf2, 'Разработка системы классификации растений '
                       'по изображениям листьев на основе EfficientNet-B0 '
                       'с Transfer Learning и веб-интерфейсом Streamlit',
                  font_size=16, color=LIGHT_GRAY)

    # Задачи
    tf3 = add_textbox(slide, Cm(1), Cm(10.5), Cm(23), Cm(7),
                      'Задачи:', font_size=18, color=YELLOW, bold=True)
    tasks = [
        'Изучить методы компьютерного зрения и CNN',
        'Подготовить датасет PlantVillage (54 303 изображения, 38 классов)',
        'Обучить модель EfficientNet-B0 с двухэтапным Transfer Learning',
        'Разработать веб-приложение Streamlit с Grad-CAM визуализацией',
        'Провести тестирование и анализ результатов',
    ]
    for i, task in enumerate(tasks, 1):
        add_paragraph(tf3, f'  {i}. {task}', font_size=15, color=LIGHT_GRAY)


def slide_dataset(prs):
    """Слайд 3: Датасет PlantVillage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), BLUE)

    add_textbox(slide, Cm(1), Cm(0.8), Cm(23), Cm(1.5),
                'Датасет PlantVillage',
                font_size=28, color=BLACK, bold=True)

    # Характеристики
    stats = [
        ('54 303', 'изображений'),
        ('38', 'классов'),
        ('14', 'видов культур'),
        ('224x224', 'пикселей (вход)'),
    ]
    for i, (val, desc) in enumerate(stats):
        x = Cm(1 + i * 6)
        add_rounded_rect(slide, x, Cm(2.5), Cm(5), Cm(3), BG_ACCENT)
        add_textbox(slide, x, Cm(2.8), Cm(5), Cm(1.5),
                    val, font_size=36, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Cm(4.2), Cm(5), Cm(1),
                    desc, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Примеры изображений
    add_image_safe(slide, 'sample_images.png', Cm(1), Cm(6.5), width=Cm(12))
    # Распределение классов
    add_image_safe(slide, 'class_distribution.png', Cm(13.5), Cm(6.5), width=Cm(11))


def slide_architecture(prs):
    """Слайд 4: Архитектура модели."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), BLUE)

    add_textbox(slide, Cm(1), Cm(0.8), Cm(23), Cm(1.5),
                'Архитектура модели: EfficientNet-B0',
                font_size=28, color=BLACK, bold=True)

    # Схема Transfer Learning
    add_image_safe(slide, 'fig_transfer_learning.png', Cm(0.5), Cm(2.5), width=Cm(12))

    # Параметры модели
    add_rounded_rect(slide, Cm(13), Cm(2.5), Cm(11.5), Cm(8), BG_ACCENT)
    tf = add_textbox(slide, Cm(13.5), Cm(2.8), Cm(10.5), Cm(1),
                     'Параметры модели', font_size=20, color=GREEN, bold=True)

    params = [
        ('Архитектура:', 'EfficientNet-B0'),
        ('Предобучение:', 'ImageNet-1K'),
        ('Параметры:', '5.3 млн'),
        ('Dropout:', '0.3'),
        ('Оптимизатор:', 'Adam'),
        ('Этап 1 (lr):', '1e-3 (только FC)'),
        ('Этап 2 (lr):', '1e-4 (вся сеть)'),
        ('Batch size:', '32'),
        ('Loss:', 'CrossEntropy (weighted)'),
    ]
    for name, val in params:
        p = tf.add_paragraph()
        p.font.size = Pt(13)
        p.font.name = 'Calibri'
        p.space_before = Pt(4)
        run1 = p.add_run()
        run1.text = f'  {name} '
        run1.font.color.rgb = LIGHT_GRAY
        run1.font.size = Pt(15)
        run2 = p.add_run()
        run2.text = val
        run2.font.color.rgb = BLACK
        run2.font.bold = True
        run2.font.size = Pt(15)

    # Архитектура системы
    add_image_safe(slide, 'fig_system_arch.png', Cm(0.5), Cm(10.5), width=Cm(24))


def slide_training(prs):
    """Слайд 5: Обучение модели."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), BLUE)

    add_textbox(slide, Cm(1), Cm(0.8), Cm(23), Cm(1.5),
                'Процесс обучения',
                font_size=28, color=BLACK, bold=True)

    # Кривые обучения
    add_image_safe(slide, 'training_curves.png', Cm(0.5), Cm(2.5), width=Cm(15))

    # Аугментация
    add_rounded_rect(slide, Cm(16), Cm(2.5), Cm(8.5), Cm(6), BG_ACCENT)
    tf = add_textbox(slide, Cm(16.5), Cm(2.8), Cm(7.5), Cm(1),
                     'Стратегия обучения', font_size=18, color=GREEN, bold=True)
    steps = [
        'Этап 1: заморозка backbone, обучение FC (5 эпох, lr=1e-3)',
        'Этап 2: fine-tuning всей сети (15 эпох, lr=1e-4)',
        'ReduceLROnPlateau (patience=3)',
        'Early stopping (patience=7)',
        'Weighted CrossEntropy для баланса классов',
    ]
    for s in steps:
        add_paragraph(tf, f'  {s}', font_size=13, color=LIGHT_GRAY)

    # Аугментация
    add_image_safe(slide, 'augmentation_demo.png', Cm(0.5), Cm(10), width=Cm(15))

    add_rounded_rect(slide, Cm(16), Cm(10), Cm(8.5), Cm(6.5), BG_ACCENT)
    tf2 = add_textbox(slide, Cm(16.5), Cm(10.3), Cm(7.5), Cm(1),
                      'Аугментация данных', font_size=18, color=YELLOW, bold=True)
    augs = [
        'Горизонтальное отражение (p=0.5)',
        'Поворот на 15 градусов',
        'ColorJitter (яркость, контраст)',
        'Аффинные преобразования',
        'Результат: +3.1% accuracy',
    ]
    for a in augs:
        add_paragraph(tf2, f'  {a}', font_size=14, color=LIGHT_GRAY)


def slide_results(prs):
    """Слайд 6: Результаты."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), GREEN)

    add_textbox(slide, Cm(1), Cm(0.8), Cm(23), Cm(1.5),
                'Результаты классификации',
                font_size=28, color=BLACK, bold=True)

    # Метрики — большие числа
    metrics = [
        ('97.42%', 'Accuracy', GREEN),
        ('97.18%', 'Precision', BLUE),
        ('96.95%', 'Recall', ORANGE),
        ('97.04%', 'F1-score', YELLOW),
    ]
    for i, (val, name, clr) in enumerate(metrics):
        x = Cm(1 + i * 6)
        add_rounded_rect(slide, x, Cm(2.5), Cm(5.2), Cm(3.5), BG_ACCENT)
        add_textbox(slide, x, Cm(2.8), Cm(5.2), Cm(2),
                    val, font_size=36, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Cm(4.5), Cm(5.2), Cm(1),
                    name, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Матрица ошибок
    add_image_safe(slide, 'confusion_matrix.png', Cm(0.5), Cm(7), width=Cm(12))

    # Grad-CAM
    add_image_safe(slide, 'grad_cam.png', Cm(13), Cm(7), width=Cm(12))


def slide_comparison(prs):
    """Слайд 7: Сравнение моделей."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), BLUE)

    add_textbox(slide, Cm(1), Cm(0.8), Cm(23), Cm(1.5),
                'Сравнительный анализ моделей',
                font_size=28, color=BLACK, bold=True)

    # График сравнения
    add_image_safe(slide, 'model_comparison.png', Cm(0.5), Cm(2.5), width=Cm(24))

    # Таблица
    from pptx.util import Inches
    rows, cols = 5, 5
    tbl_shape = slide.shapes.add_table(rows, cols, Cm(1), Cm(11), Cm(23), Cm(5.5))
    tbl = tbl_shape.table

    headers = ['Модель', 'Accuracy', 'Параметры', 'Время (CPU)', 'FLOPs']
    data = [
        ['EfficientNet-B0', '96.2%', '5.3 млн', '45 мс', '0.39 G'],
        ['ResNet-50', '94.8%', '25.6 млн', '120 мс', '4.1 G'],
        ['MobileNet-V2', '93.1%', '3.4 млн', '35 мс', '0.3 G'],
        ['VGG-16', '92.5%', '138.4 млн', '350 мс', '15.5 G'],
    ]

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK

    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
            else:
                p.font.color.rgb = BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_ACCENT


def slide_webapp(prs):
    """Слайд 8: Веб-приложение."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), ORANGE)

    add_textbox(slide, Cm(1), Cm(0.8), Cm(23), Cm(1.5),
                'Веб-приложение Streamlit',
                font_size=28, color=BLACK, bold=True)

    # Скриншот интерфейса
    add_image_safe(slide, 'fig_streamlit_ui.png', Cm(0.5), Cm(2.5), width=Cm(14))

    # Функциональность
    add_rounded_rect(slide, Cm(15.5), Cm(2.5), Cm(9), Cm(9), BG_ACCENT)
    tf = add_textbox(slide, Cm(16), Cm(2.8), Cm(8), Cm(1),
                     'Функциональность', font_size=20, color=ORANGE, bold=True)
    features = [
        'Загрузка фото (JPG, PNG)',
        'Предсказание класса + Top-5',
        'Grad-CAM визуализация',
        'Рекомендации по лечению',
        'Работа на CPU (<1 сек)',
        'Русскоязычный интерфейс',
    ]
    for f in features:
        add_paragraph(tf, f'  {f}', font_size=16, color=LIGHT_GRAY)

    # Стек технологий
    add_rounded_rect(slide, Cm(15.5), Cm(12), Cm(9), Cm(4.5), BG_ACCENT)
    tf2 = add_textbox(slide, Cm(16), Cm(12.3), Cm(8), Cm(1),
                      'Технологический стек', font_size=18, color=YELLOW, bold=True)
    stack = ['Python 3.12', 'PyTorch + torchvision', 'Streamlit',
             'matplotlib + seaborn', 'scikit-learn']
    for s in stack:
        add_paragraph(tf2, f'  {s}', font_size=15, color=LIGHT_GRAY)

    # Пример предсказания
    add_image_safe(slide, 'single_prediction.png', Cm(0.5), Cm(10.5), width=Cm(14))


def slide_demo(prs):
    """Слайд 9: Демонстрация (placeholder)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), GREEN)

    add_textbox(slide, Cm(2), Cm(5), Cm(21), Cm(3),
                'ДЕМОНСТРАЦИЯ',
                font_size=44, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_textbox(slide, Cm(2), Cm(9), Cm(21), Cm(5),
                     'Запуск веб-приложения Streamlit',
                     font_size=24, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, 'streamlit run app.py',
                  font_size=20, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, '', font_size=10)
    add_paragraph(tf, 'Загрузка изображения  ->  Классификация  ->  Grad-CAM  ->  Рекомендации',
                  font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_conclusion(prs):
    """Слайд 10: Заключение."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), GREEN)

    add_textbox(slide, Cm(1), Cm(0.8), Cm(23), Cm(1.5),
                'Заключение',
                font_size=28, color=BLACK, bold=True)

    # Результаты
    add_rounded_rect(slide, Cm(1), Cm(2.5), Cm(11.5), Cm(10), BG_ACCENT)
    tf = add_textbox(slide, Cm(1.5), Cm(2.8), Cm(10.5), Cm(1),
                     'Основные результаты', font_size=20, color=GREEN, bold=True)
    results = [
        'Обучена модель EfficientNet-B0 с accuracy 97.42%',
        'Проведено сравнение с 4 архитектурами CNN',
        'Реализовано веб-приложение с Grad-CAM',
        'Подготовлен датасет 54 303 изображения, 38 классов',
        'Время инференса: <1 сек на CPU',
    ]
    for r in results:
        add_paragraph(tf, f'  {r}', font_size=15, color=LIGHT_GRAY)

    # Перспективы
    add_rounded_rect(slide, Cm(13), Cm(2.5), Cm(11.5), Cm(10), BG_ACCENT)
    tf2 = add_textbox(slide, Cm(13.5), Cm(2.8), Cm(10.5), Cm(1),
                      'Перспективы развития', font_size=20, color=ORANGE, bold=True)
    prospects = [
        'Расширение датасета (новые культуры)',
        'Мобильное приложение (Android/iOS)',
        'Детекция нескольких заболеваний',
        'Сегментация областей поражения',
        'Интеграция с дронами для мониторинга',
    ]
    for pr in prospects:
        add_paragraph(tf2, f'  {pr}', font_size=15, color=LIGHT_GRAY)

    # Итог
    add_rounded_rect(slide, Cm(1), Cm(13.5), Cm(23.5), Cm(3), RGBColor(0xD0, 0xD0, 0xD0))
    add_textbox(slide, Cm(2), Cm(14), Cm(21.5), Cm(2),
                'Цель работы достигнута. Все задачи выполнены.',
                font_size=22, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)


def slide_thanks(prs):
    """Слайд 11: Спасибо."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rounded_rect(slide, Cm(0), Cm(0), Cm(25.4), Cm(0.3), GREEN)

    add_textbox(slide, Cm(2), Cm(5), Cm(21), Cm(3),
                'Спасибо за внимание!',
                font_size=40, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_textbox(slide, Cm(2), Cm(9), Cm(21), Cm(3),
                     'Готов ответить на ваши вопросы',
                     font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_rounded_rect(slide, Cm(0), Cm(18.6), Cm(25.4), Cm(0.3), GREEN)


# ═══════════════════════════════════════════════════
#                      MAIN
# ═══════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    slide_title(prs)        # 1. Титульный
    slide_problem(prs)      # 2. Актуальность и цель
    slide_dataset(prs)      # 3. Датасет
    slide_architecture(prs) # 4. Архитектура
    slide_training(prs)     # 5. Обучение
    slide_results(prs)      # 6. Результаты
    slide_comparison(prs)   # 7. Сравнение моделей
    slide_webapp(prs)       # 8. Веб-приложение
    slide_demo(prs)         # 9. Демонстрация (переход к живому показу)
    slide_conclusion(prs)   # 10. Заключение
    slide_thanks(prs)       # 11. Спасибо

    filename = os.path.join(DIR, 'Presentation_VKR_v3.pptx')
    prs.save(filename)
    print(f'Presentation saved: {filename}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
